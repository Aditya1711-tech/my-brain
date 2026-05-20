import io

from app.services.documents.types import RawExtraction


def parse_pptx(file_bytes: bytes, filename: str) -> RawExtraction:
    """Extract text from a PPTX file, one slide at a time."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slide_texts: list[str] = []

    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
        if parts:
            slide_texts.append("\n".join(parts))

    return RawExtraction(
        text="\n\n---\n\n".join(slide_texts),
        page_count=len(prs.slides),
    )
